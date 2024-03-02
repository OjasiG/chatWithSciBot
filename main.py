import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter, CharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
import os
import google.generativeai as genai
from langchain.vectorstores import FAISS
import openai
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
from pptx import Presentation
from docx import Document
import pdfplumber
from io import BytesIO
import openai
import subprocess
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub
from pypdf import PdfReader
import pytesseract
from PIL import Image
import pandas as pd
import platform
import PyPDF2
import openpyxl
import speech_recognition as sr
from pydub import AudioSegment
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain



GOOGLE_API_KEY= "yourOPENAIAPIKEY"
load_dotenv()
os.getenv("GOOGLE_API_KEY")
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))



def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Set the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = r'C:\Users\Disha\Downloads\Tesseract-OCR\tesseract.exe'

# Function to perform OCR on an image data
def perform_ocr(image_data):
    img = Image.open(BytesIO(image_data))
    text = pytesseract.image_to_string(img)
    return text

def images_from_pdf(pdf_docs):
    for pdf_file in pdf_docs:
        # Read PDF and extract images
        reader = PdfReader(pdf_file)
        for page_number in range(len(reader.pages)):
            page = reader.pages[page_number]
            for img in page.images:
                # Perform OCR on image data
                extracted_text = perform_ocr(img.data)
                print("Extracted Text from Image:")
                
    return extracted_text

def png_text(png_files):
    extracted_text = ""
    for png_file in png_files:
        img = Image.open(png_file)
        # Use Tesseract to extract text from the image
        text = pytesseract.image_to_string(img)
        # Append the extracted text to the result
        extracted_text += text + "\n"  # Add a newline between texts for better readability
        return extracted_text


def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(text_chunks):
    #embeddings = OpenAIEmbeddings()
    embeddings = OpenAIEmbeddings(openai_api_key="yourOPENAIKEY")

    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore


def get_conversation_chain(vectorstore):
    #llm = ChatOpenAI()
    llm = ChatOpenAI(openai_api_key="yourOPENAIKEY")

    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})

    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain


def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}",message.content ), unsafe_allow_html=True)       #
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
            

import io

def get_text_from_text_file(txt_docs):
    text = ""
    for txt_doc in txt_docs:
        # Ensure that the uploaded file is not empty
        if txt_doc is not None:
            # Read the contents of the uploaded file as bytes
            file_contents = txt_doc.read()
            # Convert bytes to string
            file_text = file_contents.decode("utf-8")
            text += file_text
    return text

def get_text_from_latex(latex_docs):
    text = ""
    for txt_doc in latex_docs:
        # Ensure that the uploaded file is not empty
        if txt_doc is not None:
            # Read the contents of the uploaded file as bytes
            file_contents = txt_doc.read()
            # Convert bytes to string
            file_text = file_contents.decode("utf-8")
            text += file_text
    return text


import pandas as pd

def get_text_from_csv(csv_files):
    text = ""
    for csv_file in csv_files:
        # Ensure that the uploaded file is not empty
        if csv_file is not None:
            # Read the contents of the uploaded CSV file into a DataFrame
            try:
                df = pd.read_csv(csv_file)
                # Convert DataFrame to string
                csv_text = df.to_string(index=False)
                text += csv_text + "\n"  # Add a newline character between CSV files
            except Exception as e:
                print(f"Error reading CSV file '{csv_file.name}': {e}")
    return text

import pandas as pd
import os
import platform
import subprocess
import os
import platform
import subprocess
import PyPDF2         
import openpyxl

import io
import speech_recognition as sr

import os
import speech_recognition as sr
from pydub import AudioSegment


def get_text_from_excel(excel_files):
    text = ""
    for excel_file in excel_files:
        # Ensure that the uploaded file is not empty
        if excel_file is not None:
            # Load the Excel file
            workbook = openpyxl.load_workbook(excel_file)

            # Assuming there's only one sheet in the workbook, if not, you need to specify the sheet name or index
            sheet = workbook.active

            # Iterate over each row in the sheet
            for row in sheet.iter_rows(values_only=True):
                # Convert row values to a sentence
                sentence = ' '.join(str(cell) for cell in row if cell is not None)
                text += sentence + "\n"  # Add a newline character between rows

            # Close the workbook when done
            workbook.close()

    return text



def get_doc_text(doc_files):
    text = ""
    for doc_file in doc_files:
        try:
            # Read the content of the uploaded file
            doc_content = doc_file.read()
            # Run the antiword command with the file content as input
            result = subprocess.run(["antiword", "-"], input=doc_content, capture_output=True, text=True)
            text += result.stdout
        except FileNotFoundError:
            print("Error: antiword command not found. Please install antiword.")
            return ""
    return text
    

def extract_tables_from_pdf(pdf_docs):
    all_tables_text = []
    for pdf_doc in pdf_docs:
        with pdfplumber.open(pdf_doc) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    if table is not None:  # Check if table is not None
                       table_text = '\n'.join([' '.join(str(cell) for cell in row if cell is not None) for row in table])
                       all_tables_text.append(table_text)
    return all_tables_text




def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks




def get_vector_store(text_chunks):
    text_chunks = [str(chunk) for chunk in text_chunks]
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")




def get_conversational_chain():


    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n


    Answer:
    """


    model = ChatGoogleGenerativeAI(model="gemini-pro",
                             temperature=0.3)


    prompt = PromptTemplate(template = prompt_template, input_variables = ["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)


    return chain


def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model = "models/embedding-001")
   
    new_db = FAISS.load_local("faiss_index", embeddings)
    docs = new_db.similarity_search(user_question)


    chain = get_conversational_chain()

   
    response = chain(
        {"input_documents":docs, "question": user_question}
        , return_only_outputs=True)


    print(response)
    st.write("Reply: ", response["output_text"])


def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text


def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        document = Document(docx)
        for paragraph in document.paragraphs:
            text += paragraph.text + "\n"
    return text


def main():
    load_dotenv()
    st.set_page_config(page_title="SciBot",page_icon=":test_tube::dna:")
    st.write(css, unsafe_allow_html=True)
    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None
    
    st.header("Chat with SciBot :dna:")
    user_question = st.text_input("Ask me any question about your scientific documents:")
    if user_question:
        handle_userinput(user_question)

    

    with st.sidebar:
        st.subheader("Your documents\n Once uploaded,press 'Process' button below")
        
        pdf_docs = st.file_uploader("Upload your PDFs here (.pdf)", accept_multiple_files=True)
        ppt_docs = st.file_uploader("Upload your PowerPoint files here (.pptx)", accept_multiple_files=True)
        docx_docs = st.file_uploader("Upload your document files here (.docx)", accept_multiple_files=True)
        txt_docs = st.file_uploader("Upload your text files here (.txt)", accept_multiple_files=True)
        latex_docs = st.file_uploader("Upload your latex files here (.tex)", accept_multiple_files=True)
        csv_files = st.file_uploader("Upload your CSV files here (.csv)", accept_multiple_files=True)
        excel_files = st.file_uploader("Upload your Excel files here (.xlsx)", accept_multiple_files=True)
        png_files = st.file_uploader("Upload your image files here (.png)", accept_multiple_files=True)

        if st.button("Process"):
            with st.spinner("Processing"):
                text_chunks = []
                if pdf_docs:
                    raw_image=images_from_pdf(pdf_docs)
                    text_chunks.extend(get_text_chunks(raw_image))
                    raw_text = get_pdf_text(pdf_docs)
                    text_chunks.extend(get_text_chunks(raw_text))
                if ppt_docs:
                    raw_text = get_ppt_text(ppt_docs)
                    text_chunks.extend(get_text_chunks(raw_text))
                if docx_docs:
                    raw_text = get_docx_text(docx_docs)
                    text_chunks.extend(get_text_chunks(raw_text))
                if txt_docs:
                    raw_text = get_text_from_text_file(txt_docs)
                    text_chunks.extend(get_text_chunks(raw_text))
                if latex_docs:
                    raw_text = get_text_from_latex(latex_docs)
                    text_chunks.extend(get_text_chunks(raw_text))
                if csv_files:
                    raw_text = get_text_from_csv(csv_files)
                    text_chunks.extend(get_text_chunks(raw_text))
                if excel_files:
                    raw_text = get_text_from_excel(excel_files)
                    text_chunks.extend(get_text_chunks(raw_text))
                if png_files:
                    raw_text = png_text(png_files)
                    text_chunks.extend(get_text_chunks(raw_text))

                vectorstore = get_vectorstore(text_chunks)
                st.success("Done")
                st.session_state.conversation = get_conversation_chain(vectorstore)
                

if __name__ == "__main__":
    main()
